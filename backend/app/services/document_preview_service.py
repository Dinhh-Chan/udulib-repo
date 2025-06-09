import io
import logging
import tempfile
import os
import subprocess
from typing import Optional, Tuple, List
from PIL import Image, ImageDraw, ImageFont
import fitz  # PyMuPDF
from fastapi import HTTPException, status
from fastapi.responses import StreamingResponse

from app.services.minio_service import minio_service
from app.core.config import settings

# Thêm import cho việc xử lý Office files
try:
    import docx
except ImportError:
    docx = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


class DocumentPreviewService:
    """Service để xử lý preview cho tất cả loại documents"""
    
    # Các loại file được hỗ trợ
    SUPPORTED_IMAGE_TYPES = ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp']
    SUPPORTED_PDF_TYPES = ['.pdf']
    SUPPORTED_OFFICE_TYPES = ['.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx']
    SUPPORTED_TEXT_TYPES = ['.txt', '.md', '.csv']
    
    # Kích thước preview
    MAX_PREVIEW_SIZE = (800, 600)
    THUMBNAIL_SIZE = (300, 300)
    LARGE_SIZE = (1200, 900)
    
    def __init__(self):
        self.logger = logging.getLogger(__name__)
    
    def get_file_type_category(self, filename: str) -> str:
        """Phân loại file theo category"""
        if not filename:
            self.logger.error("Filename is empty or None")
            return 'unsupported'
        
        try:
            ext = minio_service.get_file_extension(filename)
            
            # Nếu không có extension thì return unsupported
            if not ext:
                self.logger.info(f"File '{filename}' has no extension, treating as unsupported")
                return 'unsupported'
                
            self.logger.debug(f"File extension for '{filename}': '{ext}'")
        except Exception as e:
            self.logger.error(f"Error getting file extension for '{filename}': {e}")
            return 'unsupported'
        if ext.lower() in self.SUPPORTED_IMAGE_TYPES:
            return 'image'
        elif ext.lower() in self.SUPPORTED_PDF_TYPES:
            return 'pdf'
        elif ext.lower() in self.SUPPORTED_OFFICE_TYPES:
            return 'office'
        elif ext.lower() in self.SUPPORTED_TEXT_TYPES:
            return 'text'
        else:
            return 'unsupported'
    
    def is_supported_file(self, filename: str) -> bool:
        """Kiểm tra file có được hỗ trợ preview không"""
        return self.get_file_type_category(filename) != 'unsupported'
    
    def is_image_file(self, filename: str) -> bool:
        """Kiểm tra xem file có phải là hình ảnh không"""
        return self.get_file_type_category(filename) == 'image'
    
    def _convert_office_to_pdf(self, input_path: str) -> Optional[str]:
        """Chuyển đổi file Office sang PDF bằng LibreOffice"""
        try:
            # Tạo thư mục tạm cho output
            output_dir = tempfile.mkdtemp()
            
            # Command để chuyển đổi
            cmd = [
                'libreoffice', '--headless', '--convert-to', 'pdf',
                '--outdir', output_dir, input_path
            ]
            
            # Chạy command với timeout
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            
            if result.returncode == 0:
                # Tìm file PDF đã tạo
                base_name = os.path.splitext(os.path.basename(input_path))[0]
                pdf_path = os.path.join(output_dir, f"{base_name}.pdf")
                if os.path.exists(pdf_path):
                    return pdf_path
            
            self.logger.error(f"LibreOffice conversion failed: {result.stderr}")
            return None
            
        except subprocess.TimeoutExpired:
            self.logger.error("LibreOffice conversion timeout")
            return None
        except Exception as e:
            self.logger.error(f"LibreOffice conversion error: {e}")
            return None
    
    def _extract_docx_content(self, file_path: str) -> str:
        """Trích xuất nội dung text từ file DOCX"""
        try:
            if docx is None:
                return "File DOCX - Cần cài đặt python-docx để đọc nội dung.\n\nBạn có thể:\n- Tải file về để mở bằng Microsoft Word\n- Chuyển đổi sang PDF để preview tốt hơn"
            
            doc = docx.Document(file_path)
            text_content = []
            
            # Đọc các paragraph
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text.strip())
            
            # Đọc nội dung từ tables nếu có
            for table in doc.tables:
                text_content.append("\n--- TABLE ---")
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_content.append(" | ".join(row_text))
                text_content.append("--- END TABLE ---\n")
            
            return "\n\n".join(text_content) if text_content else "File DOCX trống hoặc không có nội dung text."
            
        except Exception as e:
            self.logger.error(f"Error extracting DOCX content: {e}")
            return f"Lỗi khi đọc file DOCX: {str(e)}\n\nBạn có thể tải file về để mở bằng Microsoft Word."
    
    def _extract_xlsx_content(self, file_path: str) -> str:
        """Trích xuất nội dung từ file XLSX"""
        try:
            if openpyxl is None:
                return "File XLSX - Cần cài đặt openpyxl để đọc nội dung.\n\nBạn có thể:\n- Tải file về để mở bằng Microsoft Excel\n- Chuyển đổi sang PDF để preview tốt hơn"
            
            workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            text_content = []
            
            # Duyệt qua tất cả worksheets (giới hạn 3 sheets đầu)
            for sheet_name in workbook.sheetnames[:3]:
                worksheet = workbook[sheet_name]
                text_content.append(f"=== SHEET: {sheet_name} ===")
                
                # Đọc dữ liệu từ sheet (giới hạn 30 rows để tránh quá dài)
                row_count = 0
                for row in worksheet.iter_rows(max_row=30, values_only=True):
                    if row_count >= 30:
                        break
                    row_data = []
                    for cell in row:
                        if cell is not None:
                            # Giới hạn độ dài text mỗi cell
                            cell_text = str(cell)
                            if len(cell_text) > 50:
                                cell_text = cell_text[:50] + "..."
                            row_data.append(cell_text)
                    if any(cell.strip() for cell in row_data if cell.strip()):
                        text_content.append(" | ".join(row_data))
                        row_count += 1
                
                if len(workbook.sheetnames) > 3:
                    text_content.append(f"\n... và {len(workbook.sheetnames) - 3} sheet khác ...")
                text_content.append("")  # Thêm dòng trống
            
            workbook.close()
            return "\n".join(text_content) if text_content else "File Excel trống hoặc không có dữ liệu."
            
        except Exception as e:
            self.logger.error(f"Error extracting XLSX content: {e}")
            return f"Lỗi khi đọc file XLSX: {str(e)}\n\nBạn có thể tải file về để mở bằng Microsoft Excel."
    
    def _extract_pptx_content(self, file_path: str) -> str:
        """Trích xuất nội dung từ file PPTX"""
        try:
            if Presentation is None:
                return "File PPTX - Cần cài đặt python-pptx để đọc nội dung.\n\nBạn có thể:\n- Tải file về để mở bằng Microsoft PowerPoint\n- Chuyển đổi sang PDF để preview tốt hơn"
            
            prs = Presentation(file_path)
            text_content = []
            
            # Duyệt qua các slides (giới hạn 10 slides đầu)
            for i, slide in enumerate(prs.slides[:10]):
                text_content.append(f"=== SLIDE {i + 1} ===")
                
                # Đọc text từ các shapes
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    text_content.extend(slide_text)
                else:
                    text_content.append("(Slide không có nội dung text hoặc chỉ có hình ảnh)")
                
                text_content.append("")  # Thêm dòng trống
            
            if len(prs.slides) > 10:
                text_content.append(f"\n... và {len(prs.slides) - 10} slide khác ...")
            
            return "\n".join(text_content) if text_content else "File PowerPoint không có nội dung text."
            
        except Exception as e:
            self.logger.error(f"Error extracting PPTX content: {e}")
            return f"Lỗi khi đọc file PPTX: {str(e)}\n\nBạn có thể tải file về để mở bằng Microsoft PowerPoint."
    
    def _extract_office_content_fallback(self, file_path: str, file_ext: str) -> str:
        """Fallback: trích xuất nội dung từ các loại file Office khi LibreOffice không có"""
        filename = os.path.basename(file_path)
        
        if file_ext.lower() == '.docx':
            return self._extract_docx_content(file_path)
        elif file_ext.lower() == '.xlsx':
            return self._extract_xlsx_content(file_path)
        elif file_ext.lower() == '.pptx':
            return self._extract_pptx_content(file_path)
        elif file_ext.lower() in ['.doc', '.xls', '.ppt']:
            return f"File {file_ext.upper()} (định dạng cũ)\n\nĐể preview file này, bạn có thể:\n- Tải file về để mở bằng Microsoft Office\n- Chuyển đổi sang định dạng mới ({file_ext.replace('.', '.').replace('doc', 'docx').replace('xls', 'xlsx').replace('ppt', 'pptx')})\n- Xuất file dưới dạng PDF"
        else:
            return f"Định dạng file Office không được hỗ trợ: {file_ext}"
    
    def _pdf_to_images(self, pdf_path: str, max_pages: int = 5) -> List[Image.Image]:
        """Chuyển PDF thành danh sách hình ảnh"""
        try:
            doc = fitz.open(pdf_path)
            images = []
            
            # Giới hạn số trang để preview
            num_pages = min(len(doc), max_pages)
            
            for page_num in range(num_pages):
                page = doc[page_num]
                # Render với độ phân giải cao
                mat = fitz.Matrix(2.0, 2.0)  # Scale 2x
                pix = page.get_pixmap(matrix=mat)
                img_data = pix.tobytes("png")
                img = Image.open(io.BytesIO(img_data))
                images.append(img)
            
            doc.close()
            return images
            
        except Exception as e:
            self.logger.error(f"Error converting PDF to images: {e}")
            return []
    
    def _create_text_preview_image(self, text_content: str, size: Tuple[int, int] = None, is_full_preview: bool = False) -> Image.Image:
        """Tạo hình ảnh preview từ nội dung text"""
        if size is None:
            size = self.MAX_PREVIEW_SIZE
        
        try:
            # Sử dụng font mặc định
            font = ImageFont.load_default()
        except:
            font = None
        
        # Xử lý encoding cho tiếng Việt
        try:
            # Đảm bảo text_content là string UTF-8 hợp lệ
            if isinstance(text_content, bytes):
                text_content = text_content.decode('utf-8', errors='replace')
            elif not isinstance(text_content, str):
                text_content = str(text_content)
        except Exception as e:
            self.logger.warning(f"Text encoding issue: {e}")
            text_content = "Không thể hiển thị nội dung do lỗi encoding"
        
        # Chuẩn bị text dựa trên loại preview
        if is_full_preview:
            # Full preview: hiển thị nhiều nội dung hơn
            lines = text_content.split('\n')[:200]  # 200 dòng cho full preview
            max_chars = 10000  # 10k ký tự cho full preview
        else:
            # Thumbnail/normal preview: giới hạn như cũ
            lines = text_content.split('\n')[:30]  # 30 dòng
            max_chars = 1000  # 1k ký tự
        
        text_to_draw = '\n'.join(lines)
        
        if len(text_to_draw) > max_chars:
            text_to_draw = text_to_draw[:max_chars] + "..."
        
        # Xử lý Unicode cho tiếng Việt - chuyển thành ASCII safe
        import unicodedata
        try:
            # Thử normalize và loại bỏ dấu trước
            text_normalized = unicodedata.normalize('NFD', text_to_draw)
            text_no_accents = ''.join(c for c in text_normalized if unicodedata.category(c) != 'Mn')
            
            # Thay thế một số ký tự đặc biệt của tiếng Việt
            vietnamese_replacements = {
                'Đ': 'D', 'đ': 'd',
                'Ă': 'A', 'ă': 'a',
                'Â': 'A', 'â': 'a', 
                'Ê': 'E', 'ê': 'e',
                'Ô': 'O', 'ô': 'o',
                'Ơ': 'O', 'ơ': 'o',
                'Ư': 'U', 'ư': 'u',
                'Ỳ': 'Y', 'ỳ': 'y',
                'Ý': 'Y', 'ý': 'y'
            }
            
            for vn_char, ascii_char in vietnamese_replacements.items():
                text_no_accents = text_no_accents.replace(vn_char, ascii_char)
            
            # Cuối cùng, encode thành ASCII và thay thế ký tự không hỗ trợ
            text_to_draw = text_no_accents.encode('ascii', errors='replace').decode('ascii')
            
        except Exception as e:
            self.logger.warning(f"Unicode normalization failed: {e}")
            # Fallback cuối cùng: chỉ giữ ASCII cơ bản
            text_to_draw = text_to_draw.encode('ascii', errors='replace').decode('ascii')
        
        # Tính toán chiều cao cần thiết cho text
        margin = 20
        line_height = 15  # Ước tính chiều cao mỗi dòng
        estimated_lines = len(text_to_draw.split('\n'))
        estimated_height = margin * 2 + estimated_lines * line_height
        
        # Giới hạn chiều cao tối đa để tránh lỗi PIL
        MAX_HEIGHT = 60000  # Giới hạn dưới 65500 pixels
        
        # Điều chỉnh kích thước image cho full preview
        if is_full_preview and estimated_height > size[1]:
            # Tạo image cao hơn để chứa full text, nhưng không vượt quá giới hạn
            desired_height = max(estimated_height, size[1])
            final_height = min(desired_height, MAX_HEIGHT)
            final_size = (size[0], final_height)
            
            # Nếu text quá dài, cắt bớt để fit
            if estimated_height > MAX_HEIGHT:
                max_lines = (MAX_HEIGHT - margin * 2) // line_height
                lines = text_to_draw.split('\n')[:max_lines]
                text_to_draw = '\n'.join(lines) + "\n\n[... Nội dung bị cắt do quá dài ...]"
        else:
            final_size = size
        
        # Tạo image
        img = Image.new('RGB', final_size, color='white')
        draw = ImageDraw.Draw(img)
        
        # Vẽ text lên image với error handling
        try:
            draw.text((margin, margin), text_to_draw, fill='black', font=font)
        except Exception as e:
            self.logger.warning(f"Error drawing text: {e}")
            # Fallback: vẽ thông báo lỗi
            error_text = f"Lỗi hiển thị text: {str(e)[:100]}"
            draw.text((margin, margin), error_text, fill='red', font=font)
        
        return img
    
    def _resize_image(self, image: Image.Image, target_size: Tuple[int, int]) -> Image.Image:
        """Resize image giữ tỷ lệ khung hình"""
        # Tạo copy để không ảnh hưởng image gốc
        img_copy = image.copy()
        img_copy.thumbnail(target_size, Image.Resampling.LANCZOS)
        
        # Chuyển đổi về RGB nếu cần
        if img_copy.mode in ('RGBA', 'LA', 'P'):
            background = Image.new('RGB', img_copy.size, (255, 255, 255))
            if img_copy.mode == 'P':
                img_copy = img_copy.convert('RGBA')
            if img_copy.mode == 'RGBA':
                background.paste(img_copy, mask=img_copy.split()[-1])
            else:
                background.paste(img_copy)
            img_copy = background
        
        return img_copy

    def _create_multipage_preview(self, images: List[Image.Image], target_size: Optional[Tuple[int, int]] = None) -> Image.Image:
        """Tạo preview nhiều trang bằng cách ghép các trang thành 1 image dài"""
        if not images:
            raise ValueError("No images provided")
        
        # Xác định kích thước của mỗi trang
        if target_size:
            max_width = target_size[0]
        else:
            max_width = 800
        
        # Giới hạn chiều cao tối đa để tránh lỗi PIL
        MAX_HEIGHT = 60000  # Giới hạn dưới 65500 pixels
        
        # Resize tất cả images về cùng width
        resized_images = []
        current_height = 0
        spacing = 20
        
        for img in images:
            # Chuyển về RGB nếu cần
            if img.mode in ('RGBA', 'LA', 'P'):
                background = Image.new('RGB', img.size, (255, 255, 255))
                if img.mode == 'P':
                    img = img.convert('RGBA')
                if img.mode == 'RGBA':
                    background.paste(img, mask=img.split()[-1])
                else:
                    background.paste(img)
                img = background
            
            # Resize giữ tỷ lệ
            ratio = max_width / img.width
            new_height = int(img.height * ratio)
            resized_img = img.resize((max_width, new_height), Image.Resampling.LANCZOS)
            
            # Kiểm tra xem thêm ảnh này có vượt quá giới hạn không
            expected_height = current_height + new_height + (spacing if resized_images else 0)
            if expected_height > MAX_HEIGHT:
                self.logger.warning(f"Giới hạn chiều cao đã đạt ({expected_height} > {MAX_HEIGHT}), dừng thêm trang")
                break
                
            resized_images.append(resized_img)
            current_height = expected_height
        
        # Tính tổng chiều cao thực tế
        total_height = sum(img.height for img in resized_images)
        total_height += spacing * (len(resized_images) - 1)
        
        # Đảm bảo không vượt quá giới hạn
        if total_height > MAX_HEIGHT:
            total_height = MAX_HEIGHT
        
        # Tạo canvas lớn
        try:
            combined_image = Image.new('RGB', (max_width, total_height), color='white')
        except Exception as e:
            self.logger.error(f"Lỗi tạo canvas kích thước {max_width}x{total_height}: {e}")
            # Fallback: chỉ return trang đầu tiên đã resize
            return resized_images[0] if resized_images else images[0]
        
        # Dán các trang vào canvas
        y_offset = 0
        for img in resized_images:
            if y_offset + img.height > total_height:
                # Cắt ảnh nếu vượt quá chiều cao canvas
                remaining_height = total_height - y_offset
                if remaining_height > 0:
                    cropped_img = img.crop((0, 0, img.width, remaining_height))
                    combined_image.paste(cropped_img, (0, y_offset))
                break
            
            combined_image.paste(img, (0, y_offset))
            y_offset += img.height + spacing
        
        return combined_image
    
    def get_document_preview_stream(self, file_path: str, filename: str, size_type: str = "medium") -> StreamingResponse:
        """
        Lấy preview stream cho bất kỳ loại tài liệu nào được hỗ trợ
        
        Args:
            file_path: Đường dẫn file trong MinIO
            filename: Tên file để xác định loại
            size_type: Loại kích thước (small, medium, large, original, full)
        """
        try:
            file_category = self.get_file_type_category(filename)
            
            if file_category == 'unsupported':
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail="Định dạng file không được hỗ trợ preview"
                )
            
            # Xác định kích thước target và số trang
            is_full_preview = size_type == "full"
            
            if size_type == "small":
                target_size = self.THUMBNAIL_SIZE
                max_pages = 1
            elif size_type == "large":
                target_size = self.LARGE_SIZE  
                max_pages = 5
            elif size_type == "original":
                target_size = None
                max_pages = 10
            elif size_type == "full":
                target_size = self.LARGE_SIZE
                max_pages = 50  # Nhiều trang cho full preview
            else:  # medium
                target_size = self.MAX_PREVIEW_SIZE
                max_pages = 3
            
            # Xử lý theo loại file
            if file_category == 'image':
                return self._handle_image_preview(file_path, target_size)
            elif file_category == 'pdf':
                return self._handle_pdf_preview(file_path, target_size, max_pages, is_full_preview)
            elif file_category == 'office':
                return self._handle_office_preview(file_path, target_size, max_pages, is_full_preview)
            elif file_category == 'text':
                return self._handle_text_preview(file_path, target_size, is_full_preview)
                
        except HTTPException:
            raise
        except Exception as e:
            self.logger.error(f"Error getting document preview: {e}")
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail="Không thể tạo preview tài liệu"
            )
    
    def _handle_image_preview(self, file_path: str, target_size: Optional[Tuple[int, int]]) -> StreamingResponse:
        """Xử lý preview cho hình ảnh"""
        file_stream = minio_service.get_file_stream(file_path)
        if not file_stream:
            raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="File không tồn tại")
        
        image_data = file_stream.read()
        file_stream.close()
        
        if target_size is None:
            # Trả về original
            file_info = minio_service.get_file_info(file_path)
            content_type = file_info.get('content_type', 'image/jpeg') if file_info else 'image/jpeg'
            return StreamingResponse(
                io.BytesIO(image_data),
                media_type=content_type,
                headers={"Cache-Control": "public, max-age=3600", "Content-Disposition": "inline"}
            )
        
        # Resize image
        try:
            image = Image.open(io.BytesIO(image_data))
            resized_image = self._resize_image(image, target_size)
            
            output = io.BytesIO()
            resized_image.save(output, format='JPEG', quality=85, optimize=True)
            output.seek(0)
            
            return StreamingResponse(
                io.BytesIO(output.getvalue()),
                media_type="image/jpeg",
                headers={"Cache-Control": "public, max-age=3600", "Content-Disposition": "inline"}
            )
        except Exception as e:
            self.logger.error(f"Error processing image: {e}")
            # Fallback to original
            return StreamingResponse(
                io.BytesIO(image_data),
                media_type="image/jpeg",
                headers={"Content-Disposition": "inline"}
            )
    
    def _handle_pdf_preview(self, file_path: str, target_size: Optional[Tuple[int, int]], max_pages: int = 1, is_full_preview: bool = False) -> StreamingResponse:
        """Xử lý preview cho PDF"""
        file_stream = minio_service.get_file_stream(file_path)
        if not file_stream:
            raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="File không tồn tại")
        
        # Lưu PDF vào file tạm
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
            for chunk in file_stream.stream(32*1024):
                temp_file.write(chunk)
            temp_pdf_path = temp_file.name
        
        file_stream.close()
        
        try:
            # Chuyển PDF thành images
            images = self._pdf_to_images(temp_pdf_path, max_pages=max_pages)
            
            if not images:
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail="Không thể tạo preview từ PDF"
                )
            
            if is_full_preview and len(images) > 1:
                # Tạo multi-page preview - ghép các trang thành 1 image dài
                combined_image = self._create_multipage_preview(images, target_size)
                final_image = combined_image
            else:
                # Single page preview (thumbnail hoặc medium)
                final_image = images[0]
                
                if target_size:
                    final_image = self._resize_image(final_image, target_size)
            
            # Chuyển thành JPEG và return StreamingResponse
            output = io.BytesIO()
            final_image.save(output, format='JPEG', quality=85, optimize=True)
            output.seek(0)
            
            return StreamingResponse(
                io.BytesIO(output.getvalue()),
                media_type="image/jpeg",
                headers={"Cache-Control": "public, max-age=3600", "Content-Disposition": "inline"}
            )
            
        finally:
            # Xóa file tạm
            try:
                os.unlink(temp_pdf_path)
            except:
                pass
    
    def _handle_office_preview(self, file_path: str, target_size: Optional[Tuple[int, int]], max_pages: int = 1, is_full_preview: bool = False) -> StreamingResponse:
        """Xử lý preview cho Office documents"""
        file_stream = minio_service.get_file_stream(file_path)
        if not file_stream:
            raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="File không tồn tại")
        
        # Lưu file Office vào file tạm
        try:
            filename = file_path.split('/')[-1]
            file_ext = minio_service.get_file_extension(filename)
        except Exception as e:
            self.logger.error(f"Error getting file extension from path '{file_path}': {e}")
            file_ext = '.tmp'  # fallback extension
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as temp_file:
            for chunk in file_stream.stream(32*1024):
                temp_file.write(chunk)
            temp_office_path = temp_file.name
        
        file_stream.close()
        
        try:
            # Thử chuyển Office sang PDF trước
            pdf_path = self._convert_office_to_pdf(temp_office_path)
            
            if not pdf_path:
                # Fallback: đọc trực tiếp nội dung text từ Office file
                self.logger.info(f"LibreOffice không có, chuyển sang đọc trực tiếp file {file_ext}")
                
                # Trích xuất nội dung text
                text_content = self._extract_office_content_fallback(temp_office_path, file_ext)
                
                # Thêm thông tin file vào đầu nội dung
                filename = file_path.split('/')[-1]
                file_info = f"=== PREVIEW FILE: {filename} ===\n"
                file_info += f"Loai file: {file_ext.upper()}\n"
                file_info += f"Preview noi dung text (LibreOffice khong co san)\n"
                file_info += "="*60 + "\n\n"
                
                # Đảm bảo encoding UTF-8
                try:
                    full_content = file_info + text_content
                    # Test encoding
                    full_content.encode('utf-8')
                except UnicodeEncodeError as e:
                    self.logger.warning(f"Unicode encoding issue: {e}")
                    # Fallback: dùng ASCII safe strings
                    file_info_safe = f"=== PREVIEW FILE: {filename} ===\n"
                    file_info_safe += f"File type: {file_ext.upper()}\n"
                    file_info_safe += f"Text preview (no LibreOffice)\n"
                    file_info_safe += "="*60 + "\n\n"
                    full_content = file_info_safe + str(text_content)
                
                # Tạo image từ text content
                preview_size = target_size or self.MAX_PREVIEW_SIZE
                text_image = self._create_text_preview_image(full_content, preview_size, is_full_preview)
                
                # Chuyển thành JPEG và return StreamingResponse
                output = io.BytesIO()
                text_image.save(output, format='JPEG', quality=85, optimize=True)
                output.seek(0)
                
                return StreamingResponse(
                    io.BytesIO(output.getvalue()),
                    media_type="image/jpeg",
                    headers={"Cache-Control": "public, max-age=3600", "Content-Disposition": "inline"}
                )
            
            try:
                # Chuyển PDF thành images
                images = self._pdf_to_images(pdf_path, max_pages=max_pages)
                
                if not images:
                    raise HTTPException(
                        status_code=status.HTTP_400_BAD_REQUEST,
                        detail="Không thể tạo preview từ file Office"
                    )
                
                if is_full_preview and len(images) > 1:
                    # Multi-page preview cho Office documents
                    combined_image = self._create_multipage_preview(images, target_size)
                    final_image = combined_image
                else:
                    # Single page preview
                    final_image = images[0]
                    
                    if target_size:
                        final_image = self._resize_image(final_image, target_size)
                
                # Chuyển thành JPEG và return StreamingResponse
                output = io.BytesIO()
                final_image.save(output, format='JPEG', quality=85, optimize=True)
                output.seek(0)
                
                return StreamingResponse(
                    io.BytesIO(output.getvalue()),
                    media_type="image/jpeg",
                    headers={"Cache-Control": "public, max-age=3600", "Content-Disposition": "inline"}
                )
                
            finally:
                # Xóa PDF tạm
                try:
                    os.unlink(pdf_path)
                except:
                    pass
                    
        finally:
            # Xóa file Office tạm
            try:
                os.unlink(temp_office_path)
            except:
                pass
    
    def _handle_text_preview(self, file_path: str, target_size: Optional[Tuple[int, int]], is_full_preview: bool = False) -> StreamingResponse:
        """Xử lý preview cho text files"""
        file_stream = minio_service.get_file_stream(file_path)
        if not file_stream:
            raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="File không tồn tại")
        
        # Đọc nội dung text
        text_data = file_stream.read()
        file_stream.close()
        
        try:
            # Decode text (thử các encoding phổ biến)
            text_content = None
            for encoding in ['utf-8', 'utf-16', 'latin1', 'cp1252']:
                try:
                    text_content = text_data.decode(encoding)
                    break
                except:
                    continue
            
            if text_content is None:
                text_content = "Không thể đọc nội dung file text"
            
            # Tạo image từ text
            preview_size = target_size or self.MAX_PREVIEW_SIZE
            text_image = self._create_text_preview_image(text_content, preview_size, is_full_preview)
            
            output = io.BytesIO()
            text_image.save(output, format='JPEG', quality=85, optimize=True)
            output.seek(0)
            
            return StreamingResponse(
                io.BytesIO(output.getvalue()),
                media_type="image/jpeg",
                headers={"Cache-Control": "public, max-age=3600", "Content-Disposition": "inline"}
            )
            
        except Exception as e:
            self.logger.error(f"Error processing text file: {e}")
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Không thể tạo preview từ file text"
            )
    
    # Backwards compatibility methods
    def get_image_preview_stream(self, file_path: str, max_size: Optional[Tuple[int, int]] = None) -> StreamingResponse:
        """Backwards compatibility cho image preview"""
        filename = file_path.split('/')[-1]
        size_type = "medium"
        if max_size == self.THUMBNAIL_SIZE:
            size_type = "small"
        elif max_size == self.LARGE_SIZE:
            size_type = "large"
        return self.get_document_preview_stream(file_path, filename, size_type)
    
    def get_thumbnail_stream(self, file_path: str) -> StreamingResponse:
        """Backwards compatibility cho thumbnail"""
        filename = file_path.split('/')[-1]
        return self.get_document_preview_stream(file_path, filename, "small")
    
    def get_original_image_stream(self, file_path: str) -> StreamingResponse:
        """Backwards compatibility cho original image"""
        filename = file_path.split('/')[-1]
        return self.get_document_preview_stream(file_path, filename, "original")


document_preview_service = DocumentPreviewService() 